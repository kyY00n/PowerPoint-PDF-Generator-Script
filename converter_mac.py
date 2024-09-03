import os
from pptx import Presentation
import pandas as pd
from colorama import init, Fore, Style
import subprocess

init()

def replace_text_in_slide(slide, old_text, new_text):
    # Replaces old_text with new_text in a slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            if old_text in shape.text:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

def create_pdf_from_pptx(ppt_filename, pdf_filename):
    # Use LibreOffice CLI to convert PPTX to PDF
    try:
        subprocess.run([
            'soffice', 
            '--headless', 
            '--convert-to', 'pdf', 
            ppt_filename, 
            '--outdir', os.path.dirname(pdf_filename)
        ], check=True)
        name = pdf_filename.split("/")[-1]
        print("New file created: " + Style.BRIGHT + Fore.CYAN + name + Style.RESET_ALL)
    except subprocess.CalledProcessError as e:
        print(Fore.RED + "Error converting file:" + str(e) + Style.RESET_ALL)

def capitalize_name(full_name):
    return ' '.join(word.capitalize() for word in full_name.split())

def process_names(names_list, pptx_file_path):
    # Process a list of names, creating a PDF for each using the appropriate template based on name length
    output_folder = "output_pdfs"

    # Create the output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for name in names_list:
        # Load the PowerPoint template
        prs = Presentation(pptx_file_path)
        slide = prs.slides[0]
        replace_text_in_slide(slide, "[NAME]", capitalize_name(name))

        temp_ppt_path = os.path.join(output_folder, f"{name}.pptx")
        pdf_path = os.path.join(output_folder, f"{name}.pdf")
        prs.save(temp_ppt_path)

        # Convert to PDF
        create_pdf_from_pptx(temp_ppt_path, pdf_path)

        # Remove the temporary PowerPoint file
        os.remove(temp_ppt_path)

# Example runner code (assuming this is how you want to run it)
if __name__ == "__main__":
    print(Style.BRIGHT + Fore.GREEN + "Processing names..." + Style.RESET_ALL)
    
    # Sample CSV and PPTX file path (replace with your actual file paths)
    # csv_file_path = "aws-cloud-club-at-ewha-womans-university_answers_1701262319.csv"  # Replace with the actual CSV file path
    # pptx_file_path = "long name.pptx"  # Replace with the actual PPTX template file path

    print("Welcome to the Name to PDF Converter!" + Style.RESET_ALL)
            
    # Get the current working directory
    current_directory = os.getcwd()
            
    # List all files in the current directory
    all_files = os.listdir(current_directory)

    # Filter and list only the CSV files
    csv_files = [file for file in all_files if file.endswith('.csv')]

    pptx_files = [file for file in all_files if file.endswith('.pptx')]

    print(Style.BRIGHT + Fore.BLUE+ "\nCSV files found in current directory:" + Style.RESET_ALL)

    # Print the list of CSV files
    for i in range(len(csv_files)):
        print(Style.BRIGHT+  f"{i}. {csv_files[i]}")
        

    index1 = int(input(Style.BRIGHT + Fore.BLUE+"\nPlease enter the index for the CSV file:\n"+ Style.RESET_ALL + "eg. 0 (to choose Cloud_101_with_AWS.csv)\n-->"))
    csv_file_path = os.path.join(os.getcwd(), csv_files[index1])

    # Check if the file exists
    if not os.path.exists(csv_file_path):
        print(f"File not found: {csv_file_path}")
        exit()



    print(Style.BRIGHT + Fore.BLUE+"\nPPTX files found in current directory:"+ Style.RESET_ALL)

    # Print the list of CSV files
    for i in range(len(pptx_files)):
        print(Style.BRIGHT+ f"{i}. {pptx_files[i]}"+ Style.RESET_ALL)
        

    index2 = int(input(Style.BRIGHT + Fore.BLUE+"\nPlease enter the index for the PPTX file:\n"+ Style.RESET_ALL + "eg. 0 (to choose name.pptx)\n-->"))
    pptx_file_path = pptx_files[index2]
    # Read names from CSV
    data = pd.read_csv(csv_file_path)
    temp_names = data['Name'].tolist()
    names = [name.strip() for name in temp_names]

    process_names(names, pptx_file_path)